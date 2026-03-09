# Extract text from PDF files

import logging
import shutil
from io import BytesIO
from PyPDF2 import PdfReader

try:
    from pdf2image import convert_from_bytes, convert_from_path
    import pytesseract
    OCR_LIBRARIES_AVAILABLE = True
except ImportError:
    OCR_LIBRARIES_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

log = logging.getLogger(__name__)


def _ocr_binaries_available():
    return shutil.which("tesseract") and shutil.which("pdftoppm")


def _extract_text_with_ocr(pdf):
    if not OCR_LIBRARIES_AVAILABLE:
        raise ValueError("OCR Python packages are not available in this environment.")

    if not _ocr_binaries_available():
        raise ValueError(
            "This PDF appears to be image-based. OCR requires system tools that are not available here: tesseract and pdftoppm (Poppler)."
        )

    if isinstance(pdf, BytesIO):
        pdf.seek(0)
        images = convert_from_bytes(pdf.read())
    else:
        images = convert_from_path(pdf)

    ocr_text_parts = []
    for page_number, image in enumerate(images, start=1):
        page_text = pytesseract.image_to_string(image).strip()
        log.info(f"OCR page {page_number}: {len(page_text)} chars")
        if page_text:
            ocr_text_parts.append(page_text)

    ocr_text = "\n".join(ocr_text_parts).strip()
    if not ocr_text:
        raise ValueError("OCR ran but no readable text was found in the PDF.")

    log.info(f"Extracted {len(ocr_text)} chars via OCR")
    return ocr_text


def extract_text_from_pdf(pdf):
    """Extract text from PDFs, using OCR as a fallback for image-based files."""
    try:
        if isinstance(pdf, BytesIO):
            reader = PdfReader(pdf)
            pdf.seek(0)
        else:
            reader = PdfReader(pdf)
        
        txt = ""
        for p in reader.pages:
            pt = p.extract_text()
            if pt:
                txt += pt + "\n"
        
        txt = txt.strip()

        if txt:
            log.info(f"Extracted {len(txt)} chars via PyPDF2")
            return txt

        log.info("No selectable text found. Trying OCR fallback.")
        return _extract_text_with_ocr(pdf)
        
    except Exception as e:
        log.error(f"Extract error: {str(e)}")
        raise ValueError(f"Failed to extract: {str(e)}")


def extract_text_from_pptx(pptx):
    try:
        if not PPTX_AVAILABLE:
            raise ValueError("python-pptx not installed")
        
        if isinstance(pptx, BytesIO):
            prs = Presentation(pptx)
        else:
            prs = Presentation(pptx)
        
        txt = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt += shape.text + "\n"
        
        txt = txt.strip()
        log.info(f"Extracted {len(txt)} chars from PPTX")
        return txt
        
    except Exception as e:
        log.error(f"PPTX extract error: {str(e)}")
        raise ValueError(f"Failed to extract from PPTX: {str(e)}")


def extract_text_from_txt(txt_file):
    try:
        if isinstance(txt_file, BytesIO):
            content = txt_file.read().decode('utf-8', errors='ignore')
        else:
            with open(txt_file, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        
        content = content.strip()
        log.info(f"Extracted {len(content)} chars from TXT")
        return content
        
    except Exception as e:
        log.error(f"TXT extract error: {str(e)}")
        raise ValueError(f"Failed to extract from TXT: {str(e)}")
