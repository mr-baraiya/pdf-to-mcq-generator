# Extract text from PPTX files

import logging
from io import BytesIO
from pptx import Presentation

log = logging.getLogger(__name__)


def extract_text_from_pptx(pptx):
    """Extract text from PPTX"""
    try:
        # Open PPTX
        if isinstance(pptx, BytesIO):
            prs = Presentation(pptx)
        else:
            # From file path
            with open(pptx, 'rb') as f:
                prs = Presentation(f)
        
        # Get text from all slides
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    # iterate over runs to get text
                    # or just use paragraph.text
                    text_runs.append(paragraph.text)
        
        txt = "\n".join(text_runs)
        
        log.info(f"Extracted {len(txt)} chars")
        return txt.strip()
        
    except Exception as e:
        log.error(f"Extract error: {str(e)}")
        raise ValueError(f"Failed to extract: {str(e)}")
